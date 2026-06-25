from pptx import Presentation
import io

def extract_text_from_ppt(file_bytes: bytes) -> str:
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        raise ValueError(f"Failed to extract PowerPoint text: {e}")